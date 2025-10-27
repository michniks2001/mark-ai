"""Pitch deck and script generation from repository analysis."""
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import json


# Predefined target audiences with specific focus areas
TARGET_AUDIENCES = {
    "seed_investors": {
        "label": "Seed Stage Investors",
        "focus": "problem-solution fit, market opportunity, founding team, early traction, funding ask",
        "tone": "compelling, data-driven, visionary"
    },
    "series_a_investors": {
        "label": "Series A Investors",
        "focus": "product-market fit, growth metrics, unit economics, competitive advantage, scaling strategy",
        "tone": "metrics-focused, strategic, proven"
    },
    "enterprise_buyers": {
        "label": "Enterprise Buyers",
        "focus": "ROI, security, integration, support, compliance, case studies",
        "tone": "professional, trustworthy, technical"
    },
    "technical_team": {
        "label": "Technical Team / Engineers",
        "focus": "architecture, tech stack, scalability, code quality, developer experience, documentation",
        "tone": "technical, detailed, pragmatic",
        "market_focus": "competitive differentiation, technical advantages, developer ecosystem, integration capabilities",
        "competitor_angle": "technical superiority, performance benchmarks, developer experience comparison"
    },
    "product_managers": {
        "label": "Product Managers",
        "focus": "features, roadmap, user experience, market fit, competitive analysis",
        "tone": "strategic, user-focused, analytical",
        "market_focus": "ease of integration, time-to-value, user adoption, feature parity",
        "competitor_angle": "feature comparison, integration ease, user experience advantages"
    },
    "general_audience": {
        "label": "General Audience",
        "focus": "what it does, why it matters, key benefits, use cases",
        "tone": "accessible, clear, engaging",
        "market_focus": "everyday use cases, practical benefits, accessibility, value proposition",
        "competitor_angle": "simplicity, usability, real-world benefits"
    },
    "seed_investors": {
        "label": "Seed Stage Investors",
        "focus": "problem-solution fit, market opportunity, founding team, early traction, funding ask",
        "tone": "compelling, data-driven, visionary",
        "market_focus": "market size, growth potential, timing, competitive landscape gaps",
        "competitor_angle": "market positioning, unique approach, competitive moats"
    },
    "series_a_investors": {
        "label": "Series A Investors",
        "focus": "product-market fit, growth metrics, unit economics, competitive advantage, scaling strategy",
        "tone": "metrics-focused, strategic, proven",
        "market_focus": "market share potential, growth metrics, competitive advantages, barriers to entry",
        "competitor_angle": "market leadership potential, sustainable advantages, winner-take-most dynamics"
    },
    "enterprise_buyers": {
        "label": "Enterprise Buyers",
        "focus": "ROI, security, integration, support, compliance, case studies",
        "tone": "professional, trustworthy, technical",
        "market_focus": "enterprise readiness, compliance, security, vendor stability",
        "competitor_angle": "enterprise features, support quality, compliance certifications"
    }
}


def get_target_audience_config(audience_key: str) -> Dict[str, str]:
    """Get configuration for a target audience."""
    return TARGET_AUDIENCES.get(audience_key, TARGET_AUDIENCES["general_audience"])


def _get_audience_specific_requirements(audience_key: str) -> str:
    """Get specific requirements for each audience type."""
    requirements = {
        "technical_team": """
- Focus on WHY this tool stands out from competitors technically
- Highlight performance benchmarks, architecture advantages, code quality
- Emphasize developer experience improvements over alternatives
- Show technical differentiators (better APIs, cleaner architecture, faster performance)
- Compare integration complexity vs competitors
- Demonstrate technical superiority with concrete examples""",
        
        "product_managers": """
- Focus on ease of integration within existing systems
- Highlight time-to-value and quick wins
- Show feature parity or advantages over competitors
- Emphasize user adoption metrics and UX improvements
- Demonstrate how it fits into current workflows
- Compare integration effort vs competitors""",
        
        "general_audience": """
- Focus on practical everyday use cases
- Explain how it makes life easier in simple terms
- Show real-world benefits anyone can understand
- Emphasize simplicity and accessibility over competitors
- Demonstrate value without technical jargon
- Compare usability vs alternatives""",
        
        "seed_investors": """
- Focus on market timing and opportunity gaps
- Highlight unique approach vs existing solutions
- Show early traction or validation
- Emphasize competitive moats and defensibility
- Demonstrate market positioning strategy
- Compare market opportunity vs alternatives""",
        
        "series_a_investors": """
- Focus on sustainable competitive advantages
- Highlight growth metrics and market share potential
- Show barriers to entry and network effects
- Emphasize winner-take-most dynamics
- Demonstrate path to market leadership
- Compare competitive positioning vs alternatives""",
        
        "enterprise_buyers": """
- Focus on enterprise readiness and compliance
- Highlight security, support, and SLAs
- Show ROI and total cost of ownership
- Emphasize vendor stability and long-term support
- Demonstrate enterprise features vs competitors
- Compare compliance certifications vs alternatives"""
    }
    
    return requirements.get(audience_key, requirements["general_audience"])


def build_pitch_deck_prompt(context: str, audience_key: str, repo_url: str) -> str:
    """Build a structured prompt for pitch deck generation."""
    audience_config = get_target_audience_config(audience_key)
    
    # Build audience-specific guidance
    market_focus = audience_config.get('market_focus', 'market opportunity and competitive landscape')
    competitor_angle = audience_config.get('competitor_angle', 'competitive advantages')
    
    return f"""You are an expert pitch deck creator and marketing strategist.

CRITICAL: Tailor this pitch deck specifically for the target audience!

Target Audience: {audience_config['label']}
Primary Focus: {audience_config['focus']}
Tone: {audience_config['tone']}
Market Analysis Focus: {market_focus}
Competitive Positioning: {competitor_angle}

Repository: {repo_url}

AUDIENCE-SPECIFIC REQUIREMENTS:
{_get_audience_specific_requirements(audience_key)}

Based on the following repository context, create a comprehensive pitch deck structure:

{context}

Generate a pitch deck with the following sections. Return ONLY valid JSON with this exact structure:

{{
  "title": "Project name or tagline",
  "slides": [
    {{
      "title": "Cover",
      "content": ["Project name", "One-line description", "Key value proposition"],
      "speaker_notes": "Opening remarks for presenter"
    }},
    {{
      "title": "Problem",
      "content": ["First key problem or pain point", "Second problem or challenge", "Third problem or market gap", "Why this matters now"],
      "speaker_notes": "How to present the problem compellingly"
    }},
    {{
      "title": "Solution",
      "content": ["Core solution approach", "How it addresses the problem", "Unique value proposition", "Key differentiator"],
      "speaker_notes": "Emphasize unique value proposition"
    }},
    {{
      "title": "Product/Features",
      "content": ["Feature 1 with brief explanation", "Feature 2 with brief explanation", "Feature 3 with brief explanation", "Additional key capability"],
      "speaker_notes": "Technical highlights and differentiators"
    }},
    {{
      "title": "Technology",
      "content": ["Primary tech stack component", "Architecture approach", "Scalability features", "Technical advantages"],
      "speaker_notes": "Why these technology choices matter"
    }},
    {{
      "title": "Market Analysis",
      "content": ["Market size and growth rate", "Key trends in the industry", "Target market segment", "Competitive positioning"],
      "speaker_notes": "Present market research and positioning strategy"
    }},
    {{
      "title": "Traction/Progress",
      "content": ["Recent development milestone", "Key improvements or features added", "Community or user engagement", "Development velocity"],
      "speaker_notes": "Demonstrate momentum and velocity"
    }},
    {{
      "title": "Roadmap",
      "content": ["Short-term goal (next 3 months)", "Medium-term goal (6-12 months)", "Long-term vision", "Key milestones"],
      "speaker_notes": "Vision for growth and expansion"
    }},
    {{
      "title": "Call to Action",
      "content": ["Primary ask or next step", "How audience can get involved", "Contact or follow-up information"],
      "speaker_notes": "Closing remarks and ask"
    }}
  ]
}}

CRITICAL REQUIREMENTS:
- Each slide MUST have AT LEAST 3 bullet points (preferably 4-5)
- Content MUST be an array of strings, NOT a single string
- Each bullet point should be informative and substantive (not just one word)
- Tailor content specifically for {audience_config['label']}
- Focus on: {audience_config['focus']}
- Use {audience_config['tone']} tone
- Make speaker notes detailed and actionable
- Return ONLY the JSON, no markdown formatting or extra text

Example of GOOD content format:
"content": [
  "Vector database using Weaviate for semantic search",
  "Next.js framework for full-stack development",
  "OpenAI GPT-4 integration for AI capabilities",
  "TypeScript for type safety and better DX"
]

Example of BAD content format (DO NOT DO THIS):
"content": "Tech stack and architecture"
"""


def create_powerpoint(pitch_data: Dict[str, Any], output_path: str) -> str:
    """
    Create a professional PowerPoint presentation from pitch deck data.
    
    Args:
        pitch_data: Structured pitch deck data with title and slides
        output_path: Path to save the .pptx file
        
    Returns:
        Path to the created file
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Style title slide
    title.text = pitch_data.get("title", "Project Pitch Deck")
    subtitle.text = "Generated from Repository Analysis"
    
    # Apply professional styling to title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)  # Professional blue
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)  # Gray
    
    # Add background color to title slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
    
    # Content slides
    for slide_data in pitch_data.get("slides", []):
        # Use title and content layout
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Style title
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "")
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)  # Professional blue
        
        # Add subtle background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        
        # Content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        content = slide_data.get("content", "")
        
        # Handle content as string or list
        if isinstance(content, list):
            lines = [str(item).strip() for item in content if str(item).strip()]
        elif isinstance(content, str):
            lines = [line.strip() for line in content.split('\n') if line.strip()]
        else:
            lines = [str(content)]
        
        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Remove bullet if already present
            clean_line = line.lstrip('•-*').strip()
            p.text = clean_line
            p.level = 0
            p.font.size = Pt(20)  # Larger, more readable
            p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray for readability
            p.space_before = Pt(12)  # Add spacing between bullets
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        speaker_notes = slide_data.get("speaker_notes", "")
        
        # Handle speaker_notes as string or list
        if isinstance(speaker_notes, list):
            text_frame.text = "\n".join(str(item) for item in speaker_notes)
        else:
            text_frame.text = str(speaker_notes)
    
    # Save
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    
    return output_path


def generate_script(pitch_data: Dict[str, Any]) -> str:
    """
    Generate a presenter script from pitch deck data.
    
    Args:
        pitch_data: Structured pitch deck data
        
    Returns:
        Formatted script as markdown
    """
    script_parts = [
        f"# Presentation Script: {pitch_data.get('title', 'Pitch Deck')}",
        "",
        "---",
        ""
    ]
    
    for i, slide_data in enumerate(pitch_data.get("slides", []), 1):
        slide_title = slide_data.get("title", f"Slide {i}")
        content = slide_data.get("content", "")
        notes = slide_data.get("speaker_notes", "")
        
        # Handle content as string or list
        if isinstance(content, list):
            content_str = "\n".join(f"- {item}" for item in content)
        else:
            content_str = str(content)
        
        # Handle notes as string or list
        if isinstance(notes, list):
            notes_str = "\n".join(str(item) for item in notes)
        else:
            notes_str = str(notes)
        
        script_parts.extend([
            f"## Slide {i}: {slide_title}",
            "",
            "**On Screen:**",
            content_str,
            "",
            "**What to Say:**",
            notes_str,
            "",
            "---",
            ""
        ])
    
    return "\n".join(script_parts)
